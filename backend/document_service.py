import os
import time
import logging
from pypdf import PdfReader
from config import settings
import csv
import openpyxl
from pptx import Presentation
import json
import pytesseract
from pdf2image import convert_from_path
import docx
import pandas as pd
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

STORE_FILE = os.path.join(settings.UPLOAD_DIR, "document_store.json")

def get_store():
    if os.path.exists(STORE_FILE):
        try:
            with open(STORE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_store(store):
    os.makedirs(settings.UPLOAD_DIR, exist_ok=True)
    with open(STORE_FILE, "w", encoding="utf-8") as f:
        json.dump(store, f)

def extract_text_from_pdf(file_path: str, filename: str) -> str:
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
            
    # Fallback to OCR if very little text is extracted (e.g. scanned PDF)
    if len(text.strip()) < 50:
        logger.info(f"Using Tesseract OCR as fallback for {filename}...")
        pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_PATH
        try:
            # Note: poppler must be in PATH or installed on the system
            images = convert_from_path(file_path)
            text = ""
            for img in images:
                text += pytesseract.image_to_string(img) + "\n"
        except Exception as e:
            logger.error(f"OCR failed for {filename}: {str(e)}")
            
    return text

def extract_text_from_csv(file_path: str, filename: str) -> str:
    try:
        df = pd.read_csv(file_path)
        # Convert to a readable string format
        return df.to_string(index=False)
    except Exception as e:
        logger.error(f"Pandas failed to read CSV {filename}: {str(e)}")
        # Fallback
        text = ""
        with open(file_path, newline='', encoding='utf-8', errors='replace') as csvfile:
            reader = csv.DictReader(csvfile)
            for row in reader:
                text += " ".join([f"{col}: {val}" for col, val in row.items() if val]) + "\n"
        return text

def extract_text_from_excel(file_path: str, filename: str) -> str:
    try:
        df = pd.read_excel(file_path, sheet_name=None)
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"--- Sheet: {sheet_name} ---\n"
            text += sheet_df.to_string(index=False) + "\n\n"
        return text
    except Exception as e:
        logger.error(f"Pandas failed to read Excel {filename}: {str(e)}")
        # Fallback
        text = ""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            headers = [cell.value for cell in sheet[1]]
            for row in sheet.iter_rows(min_row=2, values_only=True):
                row_parts = []
                for h, val in zip(headers, row):
                    if val is not None:
                        row_parts.append(f"{h}: {val}")
                text += " ".join(row_parts) + "\n"
        return text

def extract_text_from_ppt(file_path: str, filename: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file_path: str, filename: str) -> str:
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        logger.error(f"Error reading docx {filename}: {str(e)}")
    return text

def process_and_store_document(file_content: bytes, filename: str) -> int:
    os.makedirs(settings.UPLOAD_DIR, exist_ok=True)
    file_path = os.path.join(settings.UPLOAD_DIR, filename)
    with open(file_path, "wb") as f:
        f.write(file_content)
        
    ext = filename.split(".")[-1].lower()
    logger.info(f"Processing uploaded file: {filename}")
    
    text = ""
    if ext == "pdf":
        text = extract_text_from_pdf(file_path, filename)
    elif ext == "csv":
        text = extract_text_from_csv(file_path, filename)
    elif ext in ["xls", "xlsx"]:
        text = extract_text_from_excel(file_path, filename)
    elif ext in ["ppt", "pptx"]:
        text = extract_text_from_ppt(file_path, filename)
    elif ext in ["doc", "docx"]:
        text = extract_text_from_docx(file_path, filename)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
        
    if not text.strip():
        raise ValueError(f"Could not extract text from {filename}")
        
    store = get_store()
    store[filename] = text
    save_store(store)
    
    return 1

def get_all_documents() -> list:
    store = get_store()
    return list(store.keys())

def delete_document(filename: str) -> bool:
    store = get_store()
    if filename in store:
        del store[filename]
        save_store(store)
        
        file_path = os.path.join(settings.UPLOAD_DIR, filename)
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except Exception as e:
                logger.error(f"Error deleting physical file {filename}: {str(e)}")
        return True
    return False
